import pandas as pd
import os
from datetime import datetime
import cx_Oracle
import re
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN

class PM_LIST_EXPORT:
	def __init__(self, path) :
		self.ALL = None
		self.Order_Number = None
		self.SC_Number = None
		self.Customer_Code = None
		self.PATH = path
		self.ORDER_INFO()
		self.FETCH_DATA()
		self.MAKE_PM_ORDER()
		self.FILL_EXCEL_EXPORT()
		self.MAKE_PM_ORDER()

	def ORDER_INFO(self):

		for file in os.listdir(self.PATH):
			if "0210M" in file:
				path = os.path.join(self.PATH, file)
		
		self.ALL = pd.read_excel(path,dtype={'客戶產品代號(P/N)': str})
		
		#initialize the column to string 
		for col in ["華司備註", "鏈帶備註", "防鬆備註"]:
			self.ALL[col] = ""
	
		self.ALL.loc[self.ALL["穿華司"] == "Y", "華司備註"] = "有穿華司"
		self.ALL.loc[self.ALL["穿鏈帶"] == "Y", "鏈帶備註"] = "有穿鏈帶"
		self.ALL.loc[self.ALL["防鬆(電鍍後)"] == "Y", "防鬆備註"] = "有防鬆"
		self.ALL["備註"] = (
			self.ALL["華司備註"].fillna("") +
			self.ALL["鏈帶備註"].fillna("") +
			self.ALL["防鬆備註"].fillna("")
		)
		self.ALL["備註"] = self.ALL["備註"].replace("", "無")


	def FETCH_DATA(self):
		First_Item = self.ALL["客戶產品代號(P/N)"].iloc[0]
		# Define connection details
		dsn = cx_Oracle.makedsn(
			host = "192.168.1.242",      # Replace with your host IP or domain
			port = 1526,                # Replace with your port
			service_name = "sperpdb"  # Replace with your service name
		)

		# Establish the connection
		connection = cx_Oracle.connect(
			user = "spselect",         # Replace with your username
			password = "select",     # Replace with your password
			dsn = dsn
		)

		query = f"SELECT SC_NO, CST_REFE_NO, ORD_CST_NO FROM V_SCH0200Q_ORD WHERE CST_PART_NO = '{First_Item}' ORDER BY SC_NO DESC "
		df = pd.read_sql_query(query, connection)

		self.Order_Number = df["CST_REFE_NO"].iloc[0]
		self.SC_Number = df["SC_NO"].iloc[0]
		self.Customer_Code = df["ORD_CST_NO"].iloc[0]

		connection.close()

	def FILL_EXCEL_EXPORT(self):
		wb = Workbook()
		ws = wb.active

		ws["A1"] = "訂單明細"
		ws["A1"].font = Font(size=20, bold=True, underline="single")

		ws["A3"] = "客戶代號"
		ws["B3"] = self.Customer_Code
		ws["A4"] = "SC"
		ws["B4"] = self.SC_Number
		ws["A5"] = "客戶訂單號碼"
		ws["B5"] = self.Order_Number

		ws["D3"] = "訂單總箱數"
		ws["E3"] = self.ALL["訂單總箱數"].iloc[0]
		ws["D4"] = "訂單總KGS"
		ws["E4"] = self.ALL["訂單總KGS"].iloc[0]
		ws["D5"] = "訂單總M數"
		ws["E5"] = self.ALL["訂單總M數"].iloc[0]

		for row in ws.iter_rows(min_row=3, max_row=5):
			for cell in row:
				if cell.value is not None:
					cell.font = Font(size=14)

		ws["A7"] = "項次"
		ws["B7"] = "客戶產品代號(P/N)"
		ws["C7"] = "規格"
		ws["D7"] = "工程圖號"
		ws["E7"] = "訂單項次M數"
		ws["F7"] = "線徑"
		ws["G7"] = "報價千支重"
		ws["H7"] = "產品說明(中)"
		ws["I7"] = "客戶期望交期"
		ws["I7"].font = Font(bold=True, color="FF0000")
		ws["J7"] = "備註"

		header_row = 8  # Start writing data from row 8 (just below the headers)

		for i, row in self.ALL.iterrows():
			ws[f"A{header_row + i}"] = row["項次"]
			ws[f"B{header_row + i}"] = row["客戶產品代號(P/N)"]
			ws[f"C{header_row + i}"] = row["規格"]
			ws[f"D{header_row + i}"] = row["工程圖號"]
			ws[f"E{header_row + i}"] = row["訂單項次M數"]
			ws[f"F{header_row + i}"] = row["線徑"]
			ws[f"G{header_row + i}"] = row["報價千支重"]
			ws[f"H{header_row + i}"] = row["產品說明(中)"]
			ws[f"I{header_row + i}"] = row["生管交期"].strftime("%Y/%m/%d")
			ws[f"I{header_row + i}"].font = Font(bold=True, color="FF0000")
			ws[f"J{header_row + i}"] = row["備註"]


		# Define the path for saving the new Excel file
		pm_list = os.path.join(self.PATH, "生管交期確認-明細.xlsx")

		# Save the workbook to the specified path
		wb.save(pm_list)


	def MAKE_PM_ORDER(self):

		prs = Presentation()

		# 設定簡報尺寸至 A4 
		prs.slide_width = Cm(21.0)
		prs.slide_height = Cm(29.7)

		# 新增空白簡報頁面
		slide_layout = prs.slide_layouts[6]  # Blank layout
		slide = prs.slides.add_slide(slide_layout)

		content = (
			"生管訂單\n"
			f"{self.Customer_Code}\n"
			f"SC: {self.SC_Number}\n"
			f"PO: {self.Order_Number}"
		)


		# 設定文字欄大小及位置(horizontal offset from left,vertical offset from top,width,height)
		textbox = slide.shapes.add_textbox(Cm(1.0), Cm(1.0), Cm(5.5), Cm(3.5))
		text_frame = textbox.text_frame
		text_frame.clear()  # 清除其他文字欄位(若舊的存在將會有兩個content)

		p = text_frame.paragraphs[0]
		p.text = content
		p.font.size = Pt(20)
		p.alignment = PP_ALIGN.LEFT

		output_path = os.path.join(self.PATH, "生管訂單.pptx")
		prs.save(output_path)

if __name__ == "__main__":
	bot = PM_LIST_EXPORT(r"Z:\業務部\業務一課\H-訂單\1. 外銷\D09200 National Nail\1. 訂單\2025\20250609 328401(水泥板)")  # Instantiate the clas

	print("明細儲存至路徑")